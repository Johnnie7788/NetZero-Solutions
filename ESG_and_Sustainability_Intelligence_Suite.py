#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import streamlit as st
import pandas as pd
import plotly.express as px
from io import BytesIO
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

# Title and Description
st.title("ESG & Sustainability Intelligence Suite")
st.markdown("""
### ESG & Sustainability Dashboard
This tool empowers ESG analysts to efficiently manage projects, achieve regulatory compliance, and deliver impactful client solutions in the dynamic field of sustainability.
""")

# Initialize session state for datasets
if "regulatory_compliance" not in st.session_state:
    st.session_state["regulatory_compliance"] = None
if "carbon_accounting" not in st.session_state:
    st.session_state["carbon_accounting"] = None
if "esg_benchmarking" not in st.session_state:
    st.session_state["esg_benchmarking"] = None
if "project_management" not in st.session_state:
    st.session_state["project_management"] = None

# Function to load datasets
def load_uploaded_files(uploaded_file):
    try:
        return pd.read_csv(uploaded_file)
    except Exception as e:
        st.error(f"Error loading file: {e}")
        return None

# Sidebar Navigation
menu = st.sidebar.selectbox("Choose a Module", [
    "Central Dashboard",
    "Regulatory Compliance Analysis",
    "Carbon Accounting & ESG Strategy",
    "ESG Benchmarking & Initial Analysis",
    "Project Planning & Management",
    "Customer Presentation Generator",
    "Upload Required Datasets",
    "Check for Regulatory Updates"
])

# Module: Upload Required Datasets
if menu == "Upload Required Datasets":
    st.header("Upload Required Datasets")
    st.markdown("""
    Please upload the following datasets in CSV format:
    - Regulatory Compliance Data
    - Carbon Accounting Data
    - ESG Benchmarking Data
    - Project Management Data
    """)

    # Upload each dataset
    uploaded_regulatory = st.file_uploader("Upload Regulatory Compliance Dataset", type="csv")
    if uploaded_regulatory:
        st.session_state["regulatory_compliance"] = load_uploaded_files(uploaded_regulatory)
        st.success("Regulatory Compliance Dataset uploaded successfully!")

    uploaded_carbon = st.file_uploader("Upload Carbon Accounting Dataset", type="csv")
    if uploaded_carbon:
        st.session_state["carbon_accounting"] = load_uploaded_files(uploaded_carbon)
        st.success("Carbon Accounting Dataset uploaded successfully!")

    uploaded_esg = st.file_uploader("Upload ESG Benchmarking Dataset", type="csv")
    if uploaded_esg:
        st.session_state["esg_benchmarking"] = load_uploaded_files(uploaded_esg)
        st.success("ESG Benchmarking Dataset uploaded successfully!")

    uploaded_project = st.file_uploader("Upload Project Management Dataset", type="csv")
    if uploaded_project:
        st.session_state["project_management"] = load_uploaded_files(uploaded_project)
        st.success("Project Management Dataset uploaded successfully!")

# Check if datasets are loaded and not empty
if (st.session_state["regulatory_compliance"] is None or st.session_state["regulatory_compliance"].empty or
    st.session_state["carbon_accounting"] is None or st.session_state["carbon_accounting"].empty or
    st.session_state["esg_benchmarking"] is None or st.session_state["esg_benchmarking"].empty or
    st.session_state["project_management"] is None or st.session_state["project_management"].empty):
    st.sidebar.warning("Please upload all required datasets in the 'Upload Required Datasets' section.")
else:
    # Module: Central Dashboard
    if menu == "Central Dashboard":
        st.header("Central Dashboard")
        st.markdown("""
        This dashboard integrates compliance, carbon accounting, ESG benchmarking, and project planning into one cohesive view for effective project coordination.
        """)

        # Overview of Regulatory Compliance
        st.subheader("Regulatory Compliance Overview")
        st.write(st.session_state["regulatory_compliance"])

        # Overview of Carbon Accounting
        st.subheader("Carbon Accounting Overview")
        total_emissions = st.session_state["carbon_accounting"][["Scope 1 (tons CO2)", "Scope 2 (tons CO2)", "Scope 3 (tons CO2)"]].sum().sum()
        st.metric(label="Total Carbon Emissions (tons CO2)", value=total_emissions)

        # Overview of ESG Benchmarking
        st.subheader("ESG Benchmarking Overview")
        fig = px.bar(st.session_state["esg_benchmarking"], x="Company Name", y="ESG Score", title="ESG Scores by Company")
        st.plotly_chart(fig)

    # Module: Regulatory Compliance Analysis
    elif menu == "Regulatory Compliance Analysis":
        st.header("Regulatory Compliance Analysis")
        st.write(st.session_state["regulatory_compliance"])

        # Filter for specific company
        company_name = st.selectbox("Select a Company", st.session_state["regulatory_compliance"]["Company Name"].unique())
        company_data = st.session_state["regulatory_compliance"][st.session_state["regulatory_compliance"]["Company Name"] == company_name]
        st.write(f"Regulatory Compliance Details for {company_name}:")
        st.write(company_data)

    # Module: Carbon Accounting & ESG Strategy
    elif menu == "Carbon Accounting & ESG Strategy":
        st.header("Carbon Accounting & ESG Strategy")
        st.write(st.session_state["carbon_accounting"])

        # Scenario Analysis
        company_name = st.selectbox("Select a Company for Analysis", st.session_state["carbon_accounting"]["Company Name"].unique())
        company_data = st.session_state["carbon_accounting"][st.session_state["carbon_accounting"]["Company Name"] == company_name]
        if not company_data.empty:
            st.write(f"Carbon Accounting Details for {company_name}:")
            st.write(company_data)
            reduction_target = st.slider("Reduction Target (%)", 0, 100, int(company_data["Reduction Target (%)"].values[0]))
            total_emissions = company_data[["Scope 1 (tons CO2)", "Scope 2 (tons CO2)", "Scope 3 (tons CO2)"]].sum(axis=1).values[0]
            reduced_emissions = total_emissions * (1 - reduction_target / 100)
            st.metric(label="Reduced Emissions (tons CO2)", value=round(reduced_emissions, 2))

    # Module: ESG Benchmarking & Initial Analysis
    elif menu == "ESG Benchmarking & Initial Analysis":
        st.header("ESG Benchmarking & Initial Analysis")
        st.write(st.session_state["esg_benchmarking"])

        # Visualization
        metric = st.selectbox("Select a Metric to Visualize", ["ESG Score", "Carbon Intensity (CO2/$M)", "Energy Efficiency (%)", "Waste Reduction (%)"])
        fig = px.bar(st.session_state["esg_benchmarking"], x="Company Name", y=metric, title=f"{metric} by Company")
        st.plotly_chart(fig)

    # Module: Project Planning & Management
    elif menu == "Project Planning & Management":
        st.header("Project Planning & Management")
        st.write(st.session_state["project_management"])

        # Add new task
        new_task = st.text_input("New Task Description")
        new_deadline = st.date_input("Task Deadline")
        assigned_to = st.text_input("Assigned To")
        if st.button("Add Task"):
            new_entry = {"Task Description": new_task, "Deadline": new_deadline, "Status": "Not Started", "Assigned To": assigned_to}
            st.session_state["project_management"] = pd.concat([st.session_state["project_management"], pd.DataFrame([new_entry])], ignore_index=True)
            st.write("Updated Task List:")
            st.write(st.session_state["project_management"])

    # Module: Customer Presentation Generator
    elif menu == "Customer Presentation Generator":
        st.header("Customer Presentation Generator")
        project_name = st.text_input("Project Name")
        summary = st.text_area("Project Summary")

        if st.button("Generate Presentation"):
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = project_name

            content_slide = prs.slides.add_slide(prs.slide_layouts[1])
            content_slide.placeholders[1].text = summary

            output = BytesIO()
            prs.save(output)
            st.download_button(
                label="Download Presentation",
                data=output.getvalue(),
                file_name=f"{project_name}_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

    # Module: Check for Regulatory Updates
    elif menu == "Check for Regulatory Updates":
        st.header("Check for Regulatory Updates")
        st.markdown("""
        Search for the latest updates on regulatory requirements (e.g., CSRD, EU Taxonomy).
        """)

        # Input for user to search
        keyword = st.text_input("Enter a keyword to search (e.g., CSRD, EU Taxonomy)")
        if keyword:
            st.write(f"Searching for latest updates on {keyword}...")

            # Web scraping example
            url = f"https://ec.europa.eu/search/?queryText={keyword.replace(' ', '+')}"
            response = requests.get(url)

            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                results = soup.find_all('a', href=True)

                # Filter links for relevant results
                relevant_results = [
                    (result.get_text(strip=True), result['href'])
                    for result in results
                    if keyword.lower() in result.get_text(strip=True).lower() and "europa.eu" in result['href']
                ]

                # Display results or notify if none are found
                if relevant_results:
                    st.subheader("Latest Updates:")
                    for i, (text, link) in enumerate(relevant_results[:5]):  # Show top 5 results
                        st.markdown(f"{i+1}. [{text}]({link})")
                else:
                    st.warning("No relevant updates found. Try refining your search keyword.")
            else:
                st.error("Failed to fetch updates. Try again later.")

